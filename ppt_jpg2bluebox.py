import os
import shutil
import re
import threading
from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE_TYPE 
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
from tkinter import Tk, Label, Button, filedialog, messagebox

# 금지 문자 제거 및 파일명 정리
def sanitize_filename(name):
    name = re.sub(r'[\\/:"*?<>|]', '_', name)
    return name.replace('.', '_')

# 이미지 삭제 → 파란색 도형으로 대체
def replace_images_with_boxes_and_save(prs, output_path):
    for slide in prs.slides:
        for shape in list(slide.shapes):  # 복사본 순회
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height

                slide.shapes._spTree.remove(shape._element)

                # 같은 위치에 파란 사각형 도형 삽입
                rect = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height )
                fill = rect.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(0, 102, 255)  # 파란색
                rect.line.fill.background()  # 테두리 제거

    prs.save(output_path)

# 전체 실행 로직
def extract_images_from_all_pptx(input_dir, output_dir_images, output_dir_no_images):
    os.makedirs(output_dir_images, exist_ok=True)
    os.makedirs(output_dir_no_images, exist_ok=True)

    pptx_files = [f for f in os.listdir(input_dir) if f.endswith('.pptx')]
    if not pptx_files:
        messagebox.showwarning("경고", "선택한 폴더에 pptx 파일이 없습니다.")
        return

    for filename in pptx_files:
        pptx_path = os.path.join(input_dir, filename)
        prs = Presentation(pptx_path)
        image_count = 0

        ppt_name = sanitize_filename(os.path.splitext(filename)[0])
        ppt_image_folder = os.path.join(output_dir_images, ppt_name)
        os.makedirs(ppt_image_folder, exist_ok=True)

        for slide_index, slide in enumerate(prs.slides):
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    image = shape.image
                    image_bytes = image.blob
                    content_type = image.content_type
                    ext = content_type.split('/')[-1]
                    image_filename = f"slide{slide_index+1}_img{shape_index+1}.{ext}"
                    image_path = os.path.join(ppt_image_folder, image_filename)

                    try:
                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)
                        image_count += 1
                    except OSError as e:
                        print(f"[⚠️] 이미지 저장 실패: {image_path}\n{e}")

        if image_count > 0:
            output_clean_pptx = os.path.join(output_dir_no_images, filename)
            replace_images_with_boxes_and_save(prs, output_clean_pptx)
            print(f"[✅] {filename} → 이미지 {image_count}개 추출 + 파란 박스 적용")
        else:
            shutil.copy(pptx_path, os.path.join(output_dir_no_images, filename))
            shutil.rmtree(ppt_image_folder)
            print(f"[❌] {filename} → 이미지 없음 → 그대로 복사")

    messagebox.showinfo("완료", "모든 PPT 처리 완료!")

# GUI 실행
def start_extraction():
    input_dir = filedialog.askdirectory(title="PPTX 폴더 선택")
    if not input_dir:
        return

    output_dir_images = os.path.join(input_dir, 'extracted_images')
    output_dir_no_images = os.path.join(input_dir, 'no_image_pptx')

    threading.Thread(
        target=extract_images_from_all_pptx,
        args=(input_dir, output_dir_images, output_dir_no_images)
    ).start()

def run_gui():
    root = Tk()
    root.title("PPT 이미지 추출 및 파란 도형 변환기")
    root.geometry("420x200")

    Label(root, text="PPT에서 이미지를 추출하고,\n파란 도형으로 교체합니다.", pady=20).pack()
    Button(root, text="PPTX 폴더 선택 및 실행", command=start_extraction, padx=20, pady=10).pack()
    Label(root, text="© 2025 PPT Extractor", pady=20).pack()

    root.mainloop()

if __name__ == "__main__":
    run_gui()
